from pptx import Presentation
from pptx.util import Inches

filen='snowflakes.jpg'
prs = Presentation('ss.pptx')
title_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title

title.text = filen[:-4]
#subtitle.text = "python-pptx was here!"
left = Inches(2)
top = Inches(2.5)
height = Inches(4.5)
pic = slide.shapes.add_picture(filen, left, top, height=height)

prs.save('attachment.pptx')